"""一键生成 8 分钟终期汇报 PPTX（科技风 · 美化版）。

用法：
    python build_pptx.py

输出：
    docs/final_presentation.pptx

设计规范：
  - 16:9 宽屏（13.33 × 7.5 英寸）
  - 主色：深空蓝 #0F1729，强调色：电光蓝 #38BDF8，亮点：品红 #F472B6
  - 中文：微软雅黑；数字/代码：Consolas
  - 18 张幻灯片，目标总时长 8 分 10 秒
  - 严格对应课程评分标准 5 项

V2 美化要点（vs V1）：
  - 标题统一用「彩色竖条 + 标题 + 副标题 + 下划线」结构
  - 表格去掉默认网格，自定义深色表头 + 隔行交替
  - 卡片改用圆角矩形（ROUNDED_RECTANGLE），间距更舒适
  - 取消花哨的右上角 3 色方块，改为更克制的右上角细线
  - 字号层级清晰：标题 28 / 副标 14 / 数字 36-44 / 正文 13
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from lxml import etree


ROOT = Path(__file__).resolve().parent
SCREENSHOTS = ROOT / "docs" / "screenshots"
OUT = ROOT / "docs" / "final_presentation.pptx"


# =================================================================
# 主题色（v2 调整 - 更克制更精致）
# =================================================================

C_BG       = RGBColor(0x0F, 0x17, 0x29)    # 深空蓝（主背景）
C_BG_CARD  = RGBColor(0x1B, 0x25, 0x42)    # 卡片背景（比主背景亮一档）
C_BG_PANEL = RGBColor(0x14, 0x1F, 0x37)    # 表头/装饰背景
C_TITLE    = RGBColor(0xFF, 0xFF, 0xFF)    # 标题白
C_BODY     = RGBColor(0xD4, 0xDD, 0xF0)    # 正文浅蓝白
C_DIM      = RGBColor(0x7A, 0x86, 0xA8)    # 注释灰
C_ACCENT   = RGBColor(0x38, 0xBD, 0xF8)    # 电光蓝（强调）
C_HOT      = RGBColor(0xF4, 0x72, 0xB6)    # 品红（亮点/危险）
C_OK       = RGBColor(0x4A, 0xDE, 0x80)    # 青绿（成功）
C_WARN     = RGBColor(0xFB, 0xBF, 0x24)    # 琥珀（警示）
C_LINE     = RGBColor(0x2D, 0x3B, 0x5C)    # 分隔线

# 字体
F_CN       = "Microsoft YaHei"
F_CN_BOLD  = "Microsoft YaHei"
F_MONO     = "Consolas"
F_EN       = "Segoe UI"

# 画布
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# =================================================================
# 基础工具
# =================================================================

def _set_bg(slide, color=C_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, shape_type, x, y, w, h, *,
              fill=None, line=None, line_width=None):
    shp = slide.shapes.add_shape(shape_type, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_width is not None:
            shp.line.width = line_width
    return shp


def add_rect(slide, x, y, w, h, **kw):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, **kw)


def add_rounded(slide, x, y, w, h, **kw):
    shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, **kw)
    # 调整圆角半径（默认偏大）
    try:
        shp.adjustments[0] = 0.08
    except Exception:
        pass
    return shp


def add_text(slide, text, x, y, w, h, *,
             font_size=14, color=C_BODY, bold=False, italic=False,
             font_name=F_CN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_image(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        return None
    if w is not None and h is not None:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if w is not None:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    if h is not None:
        return slide.shapes.add_picture(str(path), x, y, height=h)
    return slide.shapes.add_picture(str(path), x, y)


def add_speaker_notes(slide, notes: str):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


# =================================================================
# 复合组件
# =================================================================

def add_slide_title(slide, title, subtitle=None):
    """统一的标题栏：左侧粗强调色竖条 + 标题 + 副标 + 底部细线。"""
    # 左侧强调色竖条
    add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.08), Inches(0.5),
             fill=C_HOT)
    # 主标题
    add_text(slide, title,
             Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.7),
             font_size=26, bold=True, color=C_TITLE,
             font_name=F_CN, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)
    # 副标
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.7), Inches(0.95), Inches(12.0), Inches(0.4),
                 font_size=13, color=C_DIM,
                 font_name=F_CN, align=PP_ALIGN.LEFT)
    # 底部细线
    add_rect(slide, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.02),
             fill=C_LINE)


def add_slide_number(slide, idx: int, total: int = 18):
    add_text(slide, f"{idx:02d} / {total:02d}",
             SLIDE_W - Inches(1.3), SLIDE_H - Inches(0.4),
             Inches(1.0), Inches(0.3),
             font_size=10, color=C_DIM, font_name=F_MONO,
             align=PP_ALIGN.RIGHT)
    # 左下角项目水印
    add_text(slide, "DEFI · BACKTEST PLATFORM",
             Inches(0.5), SLIDE_H - Inches(0.4),
             Inches(5.0), Inches(0.3),
             font_size=9, color=C_DIM, font_name=F_MONO)


def add_decoration_top(slide):
    """顶部装饰：渐变细线（用两段不同颜色矩形拼）。"""
    add_rect(slide, Emu(0), Emu(0), Inches(8.0), Inches(0.05), fill=C_ACCENT)
    add_rect(slide, Inches(8.0), Emu(0), Inches(5.33), Inches(0.05), fill=C_HOT)


def add_metric_card(slide, x, y, w, h, *,
                    number, label, sub=None,
                    color=C_ACCENT):
    """v2 数字卡：圆角 + 顶部色条 + 数字 + 标签。"""
    # 圆角卡片
    add_rounded(slide, x, y, w, h, fill=C_BG_CARD, line=C_LINE,
                line_width=Pt(0.5))
    # 顶部色条
    add_rect(slide, x + Inches(0.1), y + Inches(0.15),
             Inches(0.3), Inches(0.05), fill=color)
    # 大数字
    add_text(slide, str(number),
             x, y + Inches(0.35), w, Inches(1.1),
             font_size=44, color=color, bold=True,
             font_name=F_MONO, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # 标签
    add_text(slide, label,
             x, y + h - Inches(0.75), w, Inches(0.4),
             font_size=14, color=C_TITLE, bold=True,
             font_name=F_CN, align=PP_ALIGN.CENTER)
    # 副标
    if sub:
        add_text(slide, sub,
                 x + Inches(0.15), y + h - Inches(0.4), w - Inches(0.3), Inches(0.35),
                 font_size=10, color=C_DIM,
                 font_name=F_CN, align=PP_ALIGN.CENTER)


def add_table_styled(slide, headers, rows, x, y, w, h, *,
                     font_size=12, header_font_size=13,
                     header_color=C_ACCENT, body_color=C_BODY,
                     header_bg=C_BG_PANEL,
                     alt_bg=None, first_col_color=None,
                     col_widths=None, header_align=PP_ALIGN.CENTER,
                     body_align=None):
    """美化表格：自定义表头背景 + 隔行交替 + 移除默认网格。"""
    cols = len(headers)
    nrows = len(rows) + 1
    table_shape = slide.shapes.add_table(nrows, cols, x, y, w, h)
    table = table_shape.table

    # 列宽
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw / total * (w.inches if hasattr(w, "inches") else w / 914400))

    # 表头
    for j, h_text in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.08)
        cell.margin_bottom = Inches(0.08)
        p = cell.text_frame.paragraphs[0]
        p.alignment = header_align
        run = p.add_run()
        run.text = h_text
        run.font.name = F_CN
        run.font.size = Pt(header_font_size)
        run.font.bold = True
        run.font.color.rgb = header_color

    # 内容
    for i, row in enumerate(rows, start=1):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            # 隔行背景
            if alt_bg is not None and i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_BG
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.06)
            cell.margin_bottom = Inches(0.06)
            p = cell.text_frame.paragraphs[0]
            if body_align is not None:
                p.alignment = body_align
            else:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = F_CN
            run.font.size = Pt(font_size)
            if first_col_color is not None and j == 0:
                run.font.bold = True
                run.font.color.rgb = first_col_color
            else:
                run.font.color.rgb = body_color

    # 移除单元格边框（让背景色直接拼接）
    for row in table.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for ln in ('lnL', 'lnR', 'lnT', 'lnB'):
                existing = tcPr.find(qn(f"a:{ln}"))
                if existing is not None:
                    tcPr.remove(existing)
                ln_el = etree.SubElement(tcPr, qn(f"a:{ln}"))
                ln_el.set("w", "0")
                noFill = etree.SubElement(ln_el, qn("a:noFill"))

    return table_shape


def add_bullet_block(slide, items, x, y, w, h, *,
                     font_size=12, color=C_BODY,
                     bullet_color=None, bullet_char="▸",
                     line_spacing=1.4):
    """项目列表块。"""
    bc = bullet_color or C_ACCENT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run_b = p.add_run()
        run_b.text = f"{bullet_char}  "
        run_b.font.name = F_MONO
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bc
        run_b.font.bold = True
        run_t = p.add_run()
        run_t.text = item
        run_t.font.name = F_CN
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
    return tb


def section_title(slide, x, y, text, color=C_ACCENT):
    """小节标题（用强调色 mono 字体）。"""
    add_text(slide, text, x, y, Inches(6), Inches(0.4),
             font_size=12, color=color, bold=True, font_name=F_MONO)


# =================================================================
# 18 张幻灯片
# =================================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]

    # =====================================================
    # Slide 1：封面
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)

    # 顶部色带
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(0.08), fill=C_ACCENT)

    # 左侧大色块（垂直）
    add_rect(s, Emu(0), Inches(0.08), Inches(0.15), SLIDE_H - Inches(0.08),
             fill=C_HOT)

    # 左侧主内容区
    add_text(s, "FINANCIAL SOFTWARE ENGINEERING · 终期汇报",
             Inches(0.7), Inches(0.9), Inches(8.0), Inches(0.4),
             font_size=12, color=C_HOT, font_name=F_MONO, bold=True)

    add_text(s, "DeFi 收益轮动与自动复投",
             Inches(0.7), Inches(2.0), Inches(11.5), Inches(0.9),
             font_size=44, bold=True, color=C_TITLE, font_name=F_CN)
    add_text(s, "回测平台",
             Inches(0.7), Inches(2.9), Inches(11.5), Inches(0.9),
             font_size=44, bold=True, color=C_ACCENT, font_name=F_CN)

    # 副标
    add_text(s, "基于多因子评分 + 门槛约束的离线策略验证系统",
             Inches(0.7), Inches(4.05), Inches(10.5), Inches(0.5),
             font_size=15, color=C_DIM, font_name=F_CN)

    # 横分隔线
    add_rect(s, Inches(0.7), Inches(4.65), Inches(0.6), Inches(0.04),
             fill=C_HOT)

    # 关键数据小预告
    add_text(s, "13,500+ LINES   ·   325 TESTS   ·   87.5% COVERAGE   ·   REAL ON-CHAIN",
             Inches(0.7), Inches(4.85), Inches(11.5), Inches(0.4),
             font_size=11, color=C_OK, font_name=F_MONO, bold=True)

    # 右下角组员信息卡
    add_rounded(s, Inches(8.6), Inches(5.7), Inches(4.3), Inches(1.4),
                fill=C_BG_CARD, line=C_LINE, line_width=Pt(0.5))
    add_text(s, "TEAM · 2026-06-12",
             Inches(8.8), Inches(5.85), Inches(4.0), Inches(0.3),
             font_size=10, color=C_DIM, font_name=F_MONO)
    add_text(s, "金融软件工程 · 第 N 组",
             Inches(8.8), Inches(6.15), Inches(4.0), Inches(0.4),
             font_size=15, color=C_TITLE, bold=True, font_name=F_CN)
    add_text(s, "成员 A · B · C · D",
             Inches(8.8), Inches(6.55), Inches(4.0), Inches(0.4),
             font_size=12, color=C_BODY, font_name=F_CN)

    add_speaker_notes(s,
        "今天我们汇报的是 DeFi 收益轮动回测平台。一句话定位：在不接触链上的前提下，"
        "验证多池轮动 + 自动复投策略的实际收益与摩擦成本。"
    )

    # =====================================================
    # Slide 2：数字总览
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "数字总览 · CORE METRICS",
                    subtitle="代码 · 测试 · 数据 · 功能 · 性能")

    # 6 个数字卡片（2 行 3 列）
    card_w = Inches(3.95)
    card_h = Inches(2.45)
    gap = Inches(0.2)
    x0 = Inches(0.5)
    y0 = Inches(1.9)

    cards = [
        ("13,528", "代码总行数", "生产 3,747 / 测试 5,722 / UI 1,247 / 脚本 2,322", C_ACCENT),
        ("10,358", "纯代码行数", "去空白去注释 · 测试占比 41%", C_OK),
        ("325", "自动化测试", "单元 / 集成 / 黑盒 / 白盒 / 属性 / 性能", C_HOT),
        ("87.5", "% 覆盖率", "branch coverage · 核心模块 ≥ 90%", C_OK),
        ("9 + 4", "FR + NFR", "需求 100% 实现 + 5 项扩展功能", C_ACCENT),
        ("REAL", "ON-CHAIN", "DefiLlama API · 1000+ 池可选", C_HOT),
    ]
    for i, (num, lab, sub, col) in enumerate(cards):
        row = i // 3
        col_idx = i % 3
        x = x0 + col_idx * (card_w + gap)
        y = y0 + row * (card_h + gap)
        add_metric_card(s, x, y, card_w, card_h,
                        number=num, label=lab, sub=sub, color=col)

    add_slide_number(s, 2)
    add_speaker_notes(s,
        "核心数据：13500 行代码、其中纯代码 10300 多行、测试占比 41%、325 个自动化测试"
        "覆盖到 87.5% 分支、9 个功能需求全部实现并扩展了 5 项、接入了真实链上数据。"
    )

    # =====================================================
    # Slide 3：演示动画
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "看板演示 · LIVE DEMO",
                    subtitle="Streamlit 6 Tab × 10 类专业图表 · 真实链上数据")

    gif = SCREENSHOTS / "demo.gif"
    target = gif if gif.exists() else (SCREENSHOTS / "11_compare_nav_overlay.png")
    if target.exists():
        add_image(s, target, x=Inches(2.0), y=Inches(1.85),
                  w=Inches(9.3), h=Inches(5.2))

    add_text(s, "▶  放映模式下自动循环 · 10 秒一轮",
             Inches(0.5), SLIDE_H - Inches(0.55), Inches(12.3), Inches(0.3),
             font_size=11, color=C_DIM, font_name=F_CN,
             align=PP_ALIGN.CENTER)

    add_slide_number(s, 3)
    add_speaker_notes(s,
        "这是看板的实际效果。从概览到行情、风险、仓位、成本、多策略对比 6 个 Tab，"
        "每个 Tab 多张可交互图表。"
    )

    # =====================================================
    # Slide 4：实现成果矩阵
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "实现成果矩阵",
                    subtitle="9 FR + 4 NFR + 5 项扩展功能")

    headers = ["编号", "需求", "实现要点", "状态"]
    rows = [
        ["FR-01", "数据建模与 CSV 加载",   "frozen dataclass + Parquet 持久化",       "✓"],
        ["FR-02", "收益拥挤衰减模型",       "APY × TVL / (TVL + Capital)",             "✓"],
        ["FR-03", "多因子评分聚合",         "扩展为 6 个 Scorer + 权重归一化",         "✓+"],
        ["FR-04", "门槛约束型轮动",         "τ-reset + 双门槛 · 7 状态机",             "✓"],
        ["FR-05", "净效用驱动复投",         "expected_gain > gas × premium",           "✓"],
        ["FR-06", "离线回测主引擎",         "+ Mark-to-Market 重估",                   "✓+"],
        ["FR-07", "压力事件注入",           "Gas / Exploit / Liquidity 三类",          "✓"],
        ["FR-08", "绩效与归因报表",         "6 指标 + 守恒分解等式验证",               "✓"],
        ["FR-09", "交互看板",               "Streamlit 6 Tab + 多策略对比",            "✓+"],
        ["NFR",   "精度 / 复现 / 异常 / 性能", "Decimal 28 位 / 全部满足",            "✓"],
    ]
    add_table_styled(s, headers, rows,
                     Inches(0.5), Inches(1.75), Inches(7.4), Inches(5.2),
                     font_size=11, header_font_size=12,
                     alt_bg=C_BG_PANEL,
                     first_col_color=C_ACCENT,
                     col_widths=[1, 3, 5, 1])

    # 右侧扩展
    add_rounded(s, Inches(8.2), Inches(1.75), Inches(4.7), Inches(5.2),
                fill=C_BG_CARD, line=C_HOT, line_width=Pt(0.8))
    section_title(s, Inches(8.4), Inches(1.95), "超出需求的扩展", C_HOT)
    add_text(s, "EXTENSIONS",
             Inches(8.4), Inches(2.25), Inches(4.4), Inches(0.4),
             font_size=10, color=C_DIM, font_name=F_MONO)
    add_bullet_block(s, [
        "DefiLlama 真实链上数据爬取",
        "5 个命名策略预设一键切换",
        "TokenPrice 风险评分器",
        "完整测试体系（黑/白/属性/性能）",
        "GitHub Actions CI 矩阵",
    ], Inches(8.4), Inches(2.75), Inches(4.5), Inches(4.0),
       font_size=13, color=C_BODY, bullet_color=C_OK,
       line_spacing=1.6)

    add_slide_number(s, 4)
    add_speaker_notes(s,
        "9 个功能需求 + 4 个非功能需求全部实现，并且做了 5 个超出需求的扩展。"
    )

    # =====================================================
    # Slide 5：技术栈
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "技术栈 · TECHNOLOGY STACK",
                    subtitle="Python 全栈 · 零商业依赖 · 学术可复现")

    cats = [
        ("数值与精度", C_ACCENT, "NUMERIC", [
            "Python 3.13",
            "Decimal · 28 位精度",
            "NumPy 2.2 / Pandas 2.2",
        ]),
        ("持久化", C_OK, "STORAGE", [
            "PyArrow / Parquet",
            "CSV 格式（Excel 友好）",
            "YAML 配置驱动",
        ]),
        ("前端可视化", C_HOT, "FRONTEND", [
            "Streamlit 1.57",
            "Plotly 6.4",
            "Kaleido 静态导出",
        ]),
        ("链上数据", C_WARN, "ON-CHAIN", [
            "DefiLlama yields API",
            "coins.llama.fi 价格",
            "stdlib urllib（无依赖）",
        ]),
        ("测试与质量", C_ACCENT, "QUALITY", [
            "pytest + hypothesis",
            "pytest-cov branch",
            "pytest-benchmark · CI",
        ]),
    ]

    card_w_t = Inches(2.45)
    card_h_t = Inches(4.6)
    x0_t = Inches(0.5)
    y0_t = Inches(1.85)
    gap_t = Inches(0.13)

    for i, (title_t, color, en_tag, items) in enumerate(cats):
        x = x0_t + i * (card_w_t + gap_t)
        # 圆角卡片
        add_rounded(s, x, y0_t, card_w_t, card_h_t,
                    fill=C_BG_CARD, line=C_LINE, line_width=Pt(0.5))
        # 顶部色条
        add_rect(s, x, y0_t, card_w_t, Inches(0.06), fill=color)
        # EN 标签
        add_text(s, en_tag,
                 x + Inches(0.2), y0_t + Inches(0.25), card_w_t - Inches(0.4), Inches(0.3),
                 font_size=9, color=color, bold=True, font_name=F_MONO)
        # 中文标题
        add_text(s, title_t,
                 x + Inches(0.2), y0_t + Inches(0.55), card_w_t - Inches(0.4), Inches(0.5),
                 font_size=15, bold=True, color=C_TITLE, font_name=F_CN)
        # 分隔线
        add_rect(s, x + Inches(0.2), y0_t + Inches(1.15),
                 Inches(0.4), Inches(0.02), fill=color)
        # 内容
        add_bullet_block(s, items,
                         x + Inches(0.2), y0_t + Inches(1.35),
                         card_w_t - Inches(0.4), Inches(3.0),
                         font_size=11, color=C_BODY, bullet_color=color,
                         line_spacing=1.6)

    # 底部强调
    add_text(s,
             "全部金融运算用 Decimal · 链上数据无外部依赖 · 测试体系完整可执行",
             Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.4),
             font_size=12, color=C_OK, italic=True,
             font_name=F_CN, align=PP_ALIGN.CENTER)

    add_slide_number(s, 5)
    add_speaker_notes(s,
        "技术栈全部 Python 生态。两点强调：第一，所有金融计算用 Decimal 而非 float；"
        "第二，链上数据接入只用标准库 urllib，不引入新依赖。"
    )

    # =====================================================
    # Slide 6：5 层架构
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "架构设计 · 5 层模块化",
                    subtitle="依赖严格自上而下 · 抽象接口可热插拔")

    layer_specs = [
        ("ui/",         "Streamlit 交互看板 · 6 Tab + 多策略对比",   C_HOT),
        ("report/",     "metrics 6 指标 + attribution 守恒分解",     C_WARN),
        ("backtest/",   "主循环 + MTM + 摩擦 + 事件",                C_OK),
        ("strategy/",   "评分 + 轮动 + 复投 · FR-03/04/05",          C_ACCENT),
        ("data_model/", "不可变值对象 + CSV/Parquet 加载",           C_ACCENT),
    ]

    layer_x = Inches(0.6)
    layer_w = Inches(8.0)
    layer_h = Inches(0.78)
    layer_y0 = Inches(2.0)
    layer_gap = Inches(0.1)

    for i, (name, desc, color) in enumerate(layer_specs):
        y = layer_y0 + i * (layer_h + layer_gap)
        add_rounded(s, layer_x, y, layer_w, layer_h,
                    fill=C_BG_CARD, line=color, line_width=Pt(1.2))
        # 左侧色块
        add_rect(s, layer_x + Inches(0.05), y + Inches(0.12),
                 Inches(0.1), Inches(0.54), fill=color)
        add_text(s, name,
                 layer_x + Inches(0.4), y, Inches(2.5), layer_h,
                 font_size=17, bold=True, color=color,
                 font_name=F_MONO, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc,
                 layer_x + Inches(2.9), y, layer_w - Inches(3.1), layer_h,
                 font_size=13, color=C_BODY, font_name=F_CN,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 底部数据流提示
    add_text(s, "↑   依赖方向自上而下 · 低层不感知高层   ↑",
             layer_x, layer_y0 + 5 * (layer_h + layer_gap) + Inches(0.05),
             layer_w, Inches(0.4),
             font_size=11, color=C_DIM, font_name=F_CN,
             align=PP_ALIGN.CENTER)

    # 右侧设计原则
    add_rounded(s, Inches(9.0), Inches(2.0), Inches(3.85), Inches(4.7),
                fill=C_BG_CARD, line=C_LINE, line_width=Pt(0.5))
    section_title(s, Inches(9.2), Inches(2.15), "设计原则", C_HOT)
    add_bullet_block(s, [
        "frozen dataclass + tuple\n→ NFR-02 复现性",
        "Decimal 28 位精度\n→ NFR-01 金融精度",
        "IScorer 抽象接口\n→ Scorer 热插拔",
        "IFrictionEstimator 解耦\n→ 摩擦模型可替换",
        "MappingProxyType\n→ 值对象真正不可变",
    ], Inches(9.2), Inches(2.5), Inches(3.65), Inches(4.0),
       font_size=11, color=C_BODY, bullet_color=C_ACCENT,
       line_spacing=1.5)

    add_slide_number(s, 6)
    add_speaker_notes(s,
        "架构 5 层。下层是数据，往上是策略，再到回测，再到报表，最上是看板。依赖严格"
        "自上而下，每层通过抽象接口对接。"
    )

    # =====================================================
    # Slide 7：数据流
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "端到端数据流 · DATA PIPELINE",
                    subtitle="每个 tick 严格串行 · 可复现 · 365 tick / 0.19 s")

    flow_steps = [
        ("01", "数据加载",      "CSV / DefiLlama API → AssetSnapshot[]"),
        ("02", "事件注入",      "EventInjector.apply(snapshot)"),
        ("03", "MTM 重估",     "principal × (price_t / price_t-1)"),
        ("04", "收益累计",      "pending_reward += APY / 365 × principal"),
        ("05", "评分聚合",      "6 Scorer 并行 → 加权 → 稳定排序"),
        ("06", "复投判定",      "gain > gas × premium → commit"),
        ("07", "轮动判定",      "τ-reset → 双门槛 → commit"),
        ("08", "落盘报表",      "Parquet · metrics · attribution"),
    ]

    step_w = Inches(6.0)
    step_h = Inches(0.55)
    step_x = Inches(0.6)
    step_y0 = Inches(1.85)
    step_gap = Inches(0.07)

    for i, (num, name, desc) in enumerate(flow_steps):
        y = step_y0 + i * (step_h + step_gap)
        add_rounded(s, step_x, y, step_w, step_h,
                    fill=C_BG_CARD, line=C_LINE, line_width=Pt(0.5))
        # 编号小圆角
        add_rounded(s, step_x + Inches(0.1), y + Inches(0.08),
                    Inches(0.6), step_h - Inches(0.16),
                    fill=C_ACCENT)
        add_text(s, num, step_x + Inches(0.1), y + Inches(0.08),
                 Inches(0.6), step_h - Inches(0.16),
                 font_size=13, bold=True, color=C_BG,
                 font_name=F_MONO, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, name, step_x + Inches(0.85), y, Inches(1.5), step_h,
                 font_size=13, bold=True, color=C_TITLE,
                 font_name=F_CN, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc, step_x + Inches(2.35), y, step_w - Inches(2.5), step_h,
                 font_size=10, color=C_BODY, font_name=F_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 右侧亮点
    add_rounded(s, Inches(7.2), Inches(1.85), Inches(5.6), Inches(3.5),
                fill=C_BG_CARD, line=C_HOT, line_width=Pt(0.8))
    section_title(s, Inches(7.4), Inches(2.0), "关键创新 · MARK-TO-MARKET", C_HOT)
    add_text(s, "MTM 重估机制",
             Inches(7.4), Inches(2.4), Inches(5.3), Inches(0.5),
             font_size=18, bold=True, color=C_TITLE, font_name=F_CN)
    add_text(s,
             "每 tick 按 token_price 比例重估持仓\n"
             "让 NAV 反映真实价格波动\n"
             "→ MDD 不再恒为 0\n"
             "→ Sharpe 不再异常飙升",
             Inches(7.4), Inches(2.95), Inches(5.3), Inches(2.3),
             font_size=12, color=C_BODY, font_name=F_CN,
             line_spacing=1.6)

    # 性能卡
    add_rounded(s, Inches(7.2), Inches(5.5), Inches(5.6), Inches(1.5),
                fill=C_BG_CARD, line=C_OK, line_width=Pt(0.8))
    section_title(s, Inches(7.4), Inches(5.65), "性能 · PERFORMANCE", C_OK)
    add_text(s, "365 tick × 3 池  →  0.19 s",
             Inches(7.4), Inches(5.95), Inches(5.3), Inches(0.4),
             font_size=14, color=C_ACCENT, bold=True, font_name=F_MONO)
    add_text(s, "1000 tick × 5 池  →  0.75 s",
             Inches(7.4), Inches(6.35), Inches(5.3), Inches(0.4),
             font_size=14, color=C_ACCENT, bold=True, font_name=F_MONO)
    add_text(s, "NFR-04 预算 < 5 s  →  ✓ 满足",
             Inches(7.4), Inches(6.75), Inches(5.3), Inches(0.3),
             font_size=10, color=C_OK, font_name=F_CN)

    add_slide_number(s, 7)
    add_speaker_notes(s,
        "每个 tick 推进 8 步。MTM 这步是解决 Sharpe 异常的关键。性能上 1000 tick × 5 池"
        "只要 0.75 秒，远低于 5 秒预算。"
    )

    # =====================================================
    # Slide 8：真实链上数据
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "亮点 ① · 真实链上数据接入",
                    subtitle="DefiLlama 两套 API · 不引入新依赖")

    # 左：四步流程
    section_title(s, Inches(0.55), Inches(1.85), "DATA PIPELINE", C_ACCENT)

    pipeline_steps = [
        ("01", "yields.llama.fi / pools",  "列出 1000+ 池"),
        ("02", "yields.llama.fi / chart",  "单池历史 APY / TVL"),
        ("03", "coins.llama.fi / chart",   "底层 token 真实价格"),
        ("04", "归一化 + 缓存",             "日期对齐 · JSON 缓存"),
    ]
    py_y = Inches(2.3)
    for n, code, desc in pipeline_steps:
        add_rounded(s, Inches(0.55), py_y, Inches(6.0), Inches(0.85),
                    fill=C_BG_CARD, line=C_LINE, line_width=Pt(0.5))
        add_text(s, n, Inches(0.7), py_y, Inches(0.6), Inches(0.85),
                 font_size=18, bold=True, color=C_HOT,
                 font_name=F_MONO, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, code, Inches(1.3), py_y + Inches(0.1),
                 Inches(5.2), Inches(0.4),
                 font_size=13, bold=True, color=C_ACCENT, font_name=F_MONO)
        add_text(s, desc, Inches(1.3), py_y + Inches(0.45),
                 Inches(5.2), Inches(0.4),
                 font_size=11, color=C_DIM, font_name=F_CN)
        py_y += Inches(0.95)

    # 右：实测三池
    section_title(s, Inches(7.0), Inches(1.85), "实测三池 · LIVE DATA", C_ACCENT)

    headers = ["池", "APY", "token_price 范围"]
    rows = [
        ["Lido stETH",  "2.46%", "0.53 ~ 1.38"],
        ["Maple USDC",  "4.71%", "$1.00 (稳定)"],
        ["Vesper ETH",  "9.99%", "0.52 ~ 1.36"],
    ]
    add_table_styled(s, headers, rows,
                     Inches(7.0), Inches(2.3), Inches(5.85), Inches(2.3),
                     font_size=13, header_font_size=13,
                     alt_bg=C_BG_PANEL,
                     first_col_color=C_ACCENT,
                     col_widths=[2, 1.2, 2.5])

    # 关键数字
    add_rounded(s, Inches(7.0), Inches(5.0), Inches(5.85), Inches(1.7),
                fill=C_BG_CARD, line=C_OK, line_width=Pt(0.8))
    section_title(s, Inches(7.2), Inches(5.15), "覆盖范围 · COVERAGE", C_OK)
    add_text(s, "1000+ 池  ·  10+ 链  ·  300 天历史",
             Inches(7.2), Inches(5.55), Inches(5.55), Inches(0.5),
             font_size=17, color=C_TITLE, bold=True, font_name=F_MONO)
    add_text(s, "$ python fetch_data.py --demo",
             Inches(7.2), Inches(6.1), Inches(5.55), Inches(0.4),
             font_size=11, color=C_DIM, font_name=F_MONO)
    add_text(s, "一行命令拉取演示数据集",
             Inches(7.2), Inches(6.4), Inches(5.55), Inches(0.3),
             font_size=10, color=C_DIM, font_name=F_CN)

    add_slide_number(s, 8)
    add_speaker_notes(s,
        "最初只用合成数据，后来接入了 DefiLlama 的 yields 和 coins 两套 API，能拉到"
        "stETH 等真实波动 token 的历史价格。"
    )

    # =====================================================
    # Slide 9：MTM 修复
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "亮点 ② · Mark-to-Market 修复",
                    subtitle="项目最有故事的一次重构 · Sharpe 1198 → 1.5 异常修复")

    # 左：前
    add_rounded(s, Inches(0.5), Inches(1.85), Inches(6.0), Inches(2.7),
                fill=C_BG_CARD, line=C_HOT, line_width=Pt(1.2))
    section_title(s, Inches(0.7), Inches(2.0), "BEFORE · 修复前", C_HOT)
    add_text(s, "token_price ≡ 1.0 · 无 MTM",
             Inches(0.7), Inches(2.35), Inches(5.6), Inches(0.4),
             font_size=11, color=C_DIM, font_name=F_CN)
    add_text(s, "Sharpe  =  1198.376",
             Inches(0.7), Inches(2.75), Inches(5.6), Inches(0.55),
             font_size=22, bold=True, color=C_HOT, font_name=F_MONO)
    add_text(s, "MDD    =  0.00 %",
             Inches(0.7), Inches(3.3), Inches(5.6), Inches(0.55),
             font_size=22, bold=True, color=C_HOT, font_name=F_MONO)
    add_text(s, "所有策略锁定同一池 · 无法区分",
             Inches(0.7), Inches(3.95), Inches(5.6), Inches(0.4),
             font_size=12, color=C_BODY, font_name=F_CN)

    # 右：后
    add_rounded(s, Inches(6.8), Inches(1.85), Inches(6.0), Inches(2.7),
                fill=C_BG_CARD, line=C_OK, line_width=Pt(1.2))
    section_title(s, Inches(7.0), Inches(2.0), "AFTER · 修复后", C_OK)
    add_text(s, "MTM + 真实价格 + 价格风险评分",
             Inches(7.0), Inches(2.35), Inches(5.6), Inches(0.4),
             font_size=11, color=C_DIM, font_name=F_CN)
    add_text(s, "Sharpe ∈ [0.17, 1.90]",
             Inches(7.0), Inches(2.75), Inches(5.6), Inches(0.55),
             font_size=22, bold=True, color=C_OK, font_name=F_MONO)
    add_text(s, "MDD ∈ [2.15%, 21.57%]",
             Inches(7.0), Inches(3.3), Inches(5.6), Inches(0.55),
             font_size=22, bold=True, color=C_OK, font_name=F_MONO)
    add_text(s, "5 个预设产生 5 种结果 · 真正可比",
             Inches(7.0), Inches(3.95), Inches(5.6), Inches(0.4),
             font_size=12, color=C_BODY, font_name=F_CN)

    # 底部：代码
    section_title(s, Inches(0.5), Inches(4.8), "CORE LOGIC · backtest/engine.py", C_ACCENT)
    add_rounded(s, Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.85),
                fill=C_BG_PANEL, line=C_ACCENT, line_width=Pt(0.6))
    code = (
        "def _mark_to_market(self, position, prev_snap, curr_snap):\n"
        "    ratio = (curr_snap.pools[position.pool_id].token_price\n"
        "             / prev_snap.pools[position.pool_id].token_price)\n"
        "    return Position(\n"
        "        principal      = position.principal      * ratio,\n"
        "        pending_reward = position.pending_reward * ratio,\n"
        "        ...)"
    )
    add_text(s, code, Inches(0.7), Inches(5.25), Inches(12.0), Inches(1.7),
             font_size=11, color=C_BODY, font_name=F_MONO,
             align=PP_ALIGN.LEFT, line_spacing=1.3)

    add_slide_number(s, 9)
    add_speaker_notes(s,
        "项目最有故事的一次重构。最初模型没有 mark-to-market，token 涨跌完全不影响 NAV，"
        "导致 Sharpe 飙到 1198。引入 MTM 后回到正常区间。"
    )

    # =====================================================
    # Slide 10：5 预设策略
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "亮点 ③ · 5 个命名策略预设",
                    subtitle="一键切换 · 覆盖保守到激进全象限")

    # 左：参数表
    headers = ["策略", "τ", "thr", "horizon", "重心"]
    rows = [
        ["保守稳健",     "0.10", "0.002",  "90",  "价格风险 55%"],
        ["均衡（默认）", "0.05", "0.001",  "30",  "各因子均匀"],
        ["激进动量",     "0.02", "0.0005", "30",  "动量 60%"],
        ["低频价值",     "0.15", "0.005",  "120", "长 horizon"],
        ["极端风险厌恶", "0.10", "0.005",  "30",  "CARA α = 8"],
    ]
    add_table_styled(s, headers, rows,
                     Inches(0.5), Inches(1.85), Inches(7.3), Inches(3.6),
                     font_size=12, header_font_size=12,
                     alt_bg=C_BG_PANEL,
                     first_col_color=C_ACCENT,
                     col_widths=[2, 1, 1, 1.2, 2.5])

    # 右：归因雷达
    radar = SCREENSHOTS / "12_compare_radar_multi.png"
    if radar.exists():
        add_image(s, radar, x=Inches(7.95), y=Inches(1.7),
                  w=Inches(5.1), h=Inches(5.5))

    # 底部一句话
    add_rounded(s, Inches(0.5), Inches(5.7), Inches(7.3), Inches(1.3),
                fill=C_BG_CARD, line=C_OK, line_width=Pt(0.6))
    add_text(s, "🎯 文档算法原语 × 5 个权重象限 = 覆盖典型策略空间",
             Inches(0.7), Inches(5.85), Inches(7.1), Inches(0.4),
             font_size=12, color=C_OK, bold=True, font_name=F_CN)
    add_text(s, "Streamlit 侧栏一键载入 · 同步刷新所有滑块 · 切换 < 100 ms",
             Inches(0.7), Inches(6.25), Inches(7.1), Inches(0.4),
             font_size=10, color=C_DIM, font_name=F_CN)
    add_text(s, "$ # 也可手动调参，preset 设为「自定义」",
             Inches(0.7), Inches(6.55), Inches(7.1), Inches(0.4),
             font_size=10, color=C_DIM, font_name=F_MONO)

    add_slide_number(s, 10)
    add_speaker_notes(s,
        "基于文档算法原语，设计了 5 个预设覆盖保守到激进的整个象限。用户一键切换，"
        "不用手动调参。"
    )

    # =====================================================
    # Slide 11：测试体系
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "亮点 ④ · 完整测试体系",
                    subtitle="基于课程第 8/9 讲方法学 · 325 测试 / 87.5% 覆盖")

    # 左：测试金字塔
    section_title(s, Inches(0.55), Inches(1.85), "TESTING PYRAMID", C_HOT)

    pyramid = [
        ("性能基准",          "10",  C_HOT),
        ("场景 / 集成",       "21",  C_WARN),
        ("黑盒 等/边/决/场",  "88",  C_ACCENT),
        ("白盒 路径/条件",    "41",  C_OK),
        ("属性测试",          "11",  C_HOT),
        ("单元测试（模块）",  "185", C_ACCENT),
    ]
    py_y = Inches(2.3)
    for i, (name, n, color) in enumerate(pyramid):
        w_factor = 1.5 + i * 0.55
        bar_w = Inches(min(w_factor, 5.4))
        x = Inches(0.55) + (Inches(5.6) - bar_w) / 2
        add_rounded(s, x, py_y, bar_w, Inches(0.5),
                    fill=C_BG_CARD, line=color, line_width=Pt(1))
        add_text(s, f"{name}", x + Inches(0.2), py_y, bar_w - Inches(0.8), Inches(0.5),
                 font_size=11, bold=True, color=C_BODY,
                 font_name=F_CN, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, n, x + bar_w - Inches(0.6), py_y, Inches(0.5), Inches(0.5),
                 font_size=13, bold=True, color=color,
                 font_name=F_MONO, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        py_y += Inches(0.55)

    add_text(s, "Σ  325 测试",
             Inches(0.55), py_y + Inches(0.15), Inches(5.5), Inches(0.5),
             font_size=22, bold=True, color=C_HOT,
             font_name=F_MONO, align=PP_ALIGN.CENTER)

    # 右：方法学映射
    section_title(s, Inches(6.3), Inches(1.85), "PPT 方法学 · 落地映射", C_ACCENT)

    headers = ["方法（PPT §章节）", "用例", "目录"]
    rows = [
        ["等价类划分 §9.2",       "31", "blackbox/equivalence_classes"],
        ["边界值分析 §9.2",       "34", "blackbox/boundary_values"],
        ["决策表法 §9.2",         "12", "blackbox/decision_tables"],
        ["场景法 §9.2",           "11", "blackbox/scenarios"],
        ["基本路径覆盖 §一·6",   "22", "whitebox/path_coverage"],
        ["条件组合覆盖 §一·5",   "19", "whitebox/condition_combinations"],
        ["属性测试 hypothesis",  "11", "property/invariants"],
        ["性能基准 benchmark",   "10", "perf/benchmarks"],
    ]
    add_table_styled(s, headers, rows,
                     Inches(6.3), Inches(2.3), Inches(6.55), Inches(3.5),
                     font_size=10, header_font_size=11,
                     alt_bg=C_BG_PANEL,
                     first_col_color=C_ACCENT,
                     col_widths=[3.5, 1, 4])

    # 底部：测试驱动的 bug
    add_rounded(s, Inches(0.5), Inches(6.05), Inches(12.4), Inches(1.0),
                fill=C_BG_CARD, line=C_HOT, line_width=Pt(0.8))
    add_text(s, "🐛  测试驱动发现并修复 7 个真实缺陷",
             Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.45),
             font_size=14, color=C_HOT, bold=True, font_name=F_CN)
    add_text(s,
             "其中 1 个由属性测试自动发现（人想不到的越界 bug）  ·  覆盖率 87.5% · 核心策略模块 ≥ 90%",
             Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4),
             font_size=11, color=C_BODY, font_name=F_CN)

    add_slide_number(s, 11)
    add_speaker_notes(s,
        "课程第 8、9 讲讲到的等价类、边界值、决策表、路径覆盖、条件组合所有方法，我们"
        "都对应建了独立目录。325 测试发现了 7 个真实 bug，其中 1 个由属性测试自动找到。"
    )

    # =====================================================
    # Slide 12：守恒等式
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "亮点 ⑤ · 收益归因守恒等式",
                    subtitle="把「理论 - 实际」差额完整分解到 4 个可解释来源")

    # 中央等式
    add_rounded(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(1.3),
                fill=C_BG_CARD, line=C_ACCENT, line_width=Pt(1.5))
    add_text(s,
             "理论收益  =  实际收益  +  Gas  +  Slippage  +  LVR  +  调仓空窗折损",
             Inches(0.7), Inches(2.0), Inches(11.9), Inches(1.3),
             font_size=22, bold=True, color=C_ACCENT, font_name=F_CN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 4 分量卡片
    components = [
        ("Gas",      "(base + priority) × gas_limit",       C_ACCENT),
        ("Slippage", "阶梯函数 0.1 / 0.3 / 0.8 %",            C_OK),
        ("LVR",      "|oracle - pool| / oracle × size × 0.5",C_WARN),
        ("调仓空窗", "理论 - 实际 - 摩擦 = 机会成本",       C_HOT),
    ]
    col_w = Inches(2.95)
    col_x0 = Inches(0.6)
    col_y = Inches(3.6)
    for i, (name, formula, color) in enumerate(components):
        x = col_x0 + i * (col_w + Inches(0.1))
        add_rounded(s, x, col_y, col_w, Inches(1.55),
                    fill=C_BG_CARD, line=color, line_width=Pt(1.2))
        # 顶部色条
        add_rect(s, x, col_y, col_w, Inches(0.06), fill=color)
        add_text(s, name, x, col_y + Inches(0.2), col_w, Inches(0.5),
                 font_size=18, bold=True, color=color,
                 font_name=F_CN, align=PP_ALIGN.CENTER)
        add_text(s, formula, x + Inches(0.1), col_y + Inches(0.75),
                 col_w - Inches(0.2), Inches(0.7),
                 font_size=10, color=C_BODY, font_name=F_MONO,
                 align=PP_ALIGN.CENTER, line_spacing=1.5)

    # 底部：验证代码
    add_rounded(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.5),
                fill=C_BG_CARD, line=C_OK, line_width=Pt(0.8))
    section_title(s, Inches(0.8), Inches(5.6), "集成测试验证 · TEST", C_OK)
    add_text(s,
             "reconstructed = actual + gas + slippage + lvr + idle\n"
             "assert abs(reconstructed - theoretical) <= Decimal(1)",
             Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.85),
             font_size=12, color=C_BODY, font_name=F_MONO, line_spacing=1.3)
    add_text(s, "→  误差 < 1 元 / 100 k 本金  ·  即 0.001% 守恒",
             Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.3),
             font_size=10, color=C_OK, font_name=F_CN)

    add_slide_number(s, 12)
    add_speaker_notes(s,
        "把'理论最优'到'实际收益'的差额完整分解为 4 类摩擦 + 调仓空窗，并用集成测试"
        "守恒等式验证误差小于 1 元。"
    )

    # =====================================================
    # Slide 13：5 策略对比
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "演示 · 5 策略对比实测",
                    subtitle="同一份数据 · 5 种结果 · 价格风险敏感策略胜出")

    img = SCREENSHOTS / "11_compare_nav_overlay.png"
    if img.exists():
        add_image(s, img, x=Inches(0.5), y=Inches(1.85),
                  w=Inches(7.5), h=Inches(4.0))

    # 结果表
    section_title(s, Inches(8.2), Inches(1.85), "实测结果 · METRICS", C_ACCENT)
    headers = ["策略", "年化", "MDD", "Sharpe"]
    rows = [
        ["保守稳健 ✓",       "+8.82%",  "5.80%",  "1.076"],
        ["极端风险厌恶 ✓",   "+8.36%",  "5.80%",  "1.020"],
        ["均衡",             "-3.65%",  "20.69%", "-0.230"],
        ["激进动量 ✗",       "-6.50%",  "20.44%", "-0.367"],
        ["低频价值 ✗",       "-11.71%", "21.57%", "-0.781"],
    ]
    add_table_styled(s, headers, rows,
                     Inches(8.2), Inches(2.3), Inches(4.8), Inches(3.5),
                     font_size=11, header_font_size=12,
                     alt_bg=C_BG_PANEL,
                     first_col_color=C_ACCENT,
                     col_widths=[2, 1.2, 1, 1])

    # 结论
    add_rounded(s, Inches(0.5), Inches(6.05), Inches(12.4), Inches(1.0),
                fill=C_BG_CARD, line=C_OK, line_width=Pt(0.6))
    add_text(s, "✓  价格风险敏感策略避开 token 暴跌 · ✗  激进策略锁定高 APY 池被 MTM 反噬",
             Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.4),
             font_size=14, color=C_OK, bold=True, font_name=F_CN,
             align=PP_ALIGN.CENTER)
    add_text(s,
             "📊  这正是「场景设计 · 对比分析 · 结论洞察」—— 不仅能跑，还能回答问题",
             Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4),
             font_size=10, color=C_DIM, font_name=F_CN,
             align=PP_ALIGN.CENTER)

    add_slide_number(s, 13)
    add_speaker_notes(s,
        "同一份数据 5 个策略跑出完全不同的结果。价格风险敏感的策略避开 token 价格暴跌；"
        "激进策略锁定最高 APY 池被 MTM 反噬，亏损 6.5%。"
    )

    # =====================================================
    # Slide 14：真实数据 case
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "演示 · 真实链上数据 CASE STUDY",
                    subtitle="Maple USDC / Vesper ETH / Lido stETH × 300 天")

    img1 = SCREENSHOTS / "09_position_timeline.png"
    if img1.exists():
        add_image(s, img1, x=Inches(0.5), y=Inches(1.85),
                  w=Inches(8.0), h=Inches(2.4))

    img2 = SCREENSHOTS / "07_risk_drawdown.png"
    if img2.exists():
        add_image(s, img2, x=Inches(0.5), y=Inches(4.4),
                  w=Inches(8.0), h=Inches(2.4))

    # 右侧指标卡
    section_title(s, Inches(8.7), Inches(1.85), "保守稳健 · METRICS", C_ACCENT)

    metrics_data = [
        ("年化收益", "+8.74 %",  C_OK),
        ("最大回撤", "2.15 %",   C_OK),
        ("Sharpe",   "1.540",    C_ACCENT),
        ("调仓次数", "2 次",     C_BODY),
        ("USDC 占比", "99 %",    C_OK),
    ]
    my = Inches(2.3)
    for name, value, color in metrics_data:
        add_rounded(s, Inches(8.7), my, Inches(4.3), Inches(0.85),
                    fill=C_BG_CARD, line=color, line_width=Pt(0.8))
        add_text(s, name, Inches(8.9), my + Inches(0.1),
                 Inches(4.0), Inches(0.3),
                 font_size=11, color=C_DIM, font_name=F_CN)
        add_text(s, value, Inches(8.9), my + Inches(0.35),
                 Inches(4.0), Inches(0.5),
                 font_size=22, bold=True, color=color, font_name=F_MONO)
        my += Inches(0.93)

    add_slide_number(s, 14)
    add_speaker_notes(s,
        "用真实链上数据跑保守策略：年化 8.74%、MDD 2.15%、Sharpe 1.54。模型正确识别"
        "stETH 价格风险，把 99% 时间留在 Maple USDC。"
    )

    # =====================================================
    # Slide 15：成员分工
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "成员分工 · 4 人 × 25 %",
                    subtitle="按模块边界划分 · 接口先行 · 并行开发")

    members = [
        ("成员 A", "数据 + 链上",
         "data_model + onchain_fetcher",
         "CSV / Parquet 加载\nDefiLlama API 接入\nfrozen dataclass 设计\n需求文档与 README",
         C_ACCENT),
        ("成员 B", "策略引擎",
         "strategy/ Scorer + 引擎",
         "6 个 Scorer 实现\nRotationEngine 状态机\nReinvestEngine\n5 个策略预设",
         C_HOT),
        ("成员 C", "回测 + 报表",
         "backtest/ + report/",
         "FrictionEstimator\nEventInjector\nMark-to-Market 引擎\n归因守恒分解",
         C_OK),
        ("成员 D", "前端 + 测试",
         "ui/ + tests/ + CI",
         "Streamlit 6 Tab\nPlotly 10 图表\n325 测试体系\nGitHub Actions",
         C_WARN),
    ]
    card_w = Inches(3.05)
    card_h = Inches(4.6)
    x0 = Inches(0.5)
    y0 = Inches(1.85)
    gap = Inches(0.1)

    for i, (name, role, scope, contrib, color) in enumerate(members):
        x = x0 + i * (card_w + gap)
        # 圆角卡
        add_rounded(s, x, y0, card_w, card_h,
                    fill=C_BG_CARD, line=color, line_width=Pt(1.2))
        # 顶部色条
        add_rect(s, x, y0, card_w, Inches(0.08), fill=color)
        # 姓名
        add_text(s, name, x + Inches(0.2), y0 + Inches(0.25),
                 card_w - Inches(0.4), Inches(0.5),
                 font_size=20, bold=True, color=color, font_name=F_CN)
        # 角色
        add_text(s, role, x + Inches(0.2), y0 + Inches(0.75),
                 card_w - Inches(0.4), Inches(0.4),
                 font_size=13, bold=True, color=C_TITLE, font_name=F_CN)
        # 模块路径
        add_text(s, scope, x + Inches(0.2), y0 + Inches(1.15),
                 card_w - Inches(0.4), Inches(0.4),
                 font_size=9, color=C_DIM, font_name=F_MONO)
        # 分隔线
        add_rect(s, x + Inches(0.2), y0 + Inches(1.55),
                 Inches(0.4), Inches(0.02), fill=color)
        # 贡献
        add_text(s, contrib,
                 x + Inches(0.2), y0 + Inches(1.7),
                 card_w - Inches(0.4), Inches(2.2),
                 font_size=11, color=C_BODY, font_name=F_CN,
                 line_spacing=1.4)
        # 百分比
        add_text(s, "25 %", x, y0 + card_h - Inches(0.75),
                 card_w, Inches(0.6),
                 font_size=28, bold=True, color=color,
                 font_name=F_MONO, align=PP_ALIGN.CENTER)

    # 底部
    add_text(s,
             "协作机制：Git PR 互审  ·  IScorer / IFrictionEstimator 抽象接口先行  ·  每周集成测试",
             Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             font_size=11, color=C_DIM, italic=True,
             font_name=F_CN, align=PP_ALIGN.CENTER)

    add_slide_number(s, 15)
    add_speaker_notes(s,
        "4 人分工平均 25%，按模块边界划分。A 数据层和链上接入，B 策略引擎，C 回测报表，"
        "D 前端测试。接口先行让我们能并行开发。"
    )

    # =====================================================
    # Slide 16：成果与预期对比
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "成果与预期对比",
                    subtitle="✓ 多项超出预期 · ⚠ 部分难度超预期 · ↓ 少量降级")

    # 左：超出
    add_rounded(s, Inches(0.5), Inches(1.85), Inches(4.1), Inches(5.0),
                fill=C_BG_CARD, line=C_OK, line_width=Pt(1.2))
    add_rect(s, Inches(0.5), Inches(1.85), Inches(4.1), Inches(0.06), fill=C_OK)
    section_title(s, Inches(0.7), Inches(2.0), "EXCEEDED · 超出预期", C_OK)
    add_bullet_block(s, [
        "真实链上数据接入\n（原计划仅合成）",
        "5 个命名策略预设\n（原计划仅自定义）",
        "TokenPrice 风险评分器\n（课程未要求）",
        "完整测试体系\n（远超基本单测）",
        "属性测试 + CI\n（自动找到真实 bug）",
    ], Inches(0.7), Inches(2.5), Inches(3.75), Inches(4.3),
       font_size=11, color=C_BODY, bullet_color=C_OK,
       line_spacing=1.4)

    # 中：难度
    add_rounded(s, Inches(4.7), Inches(1.85), Inches(4.1), Inches(5.0),
                fill=C_BG_CARD, line=C_WARN, line_width=Pt(1.2))
    add_rect(s, Inches(4.7), Inches(1.85), Inches(4.1), Inches(0.06), fill=C_WARN)
    section_title(s, Inches(4.9), Inches(2.0), "HARD · 难度超预期", C_WARN)
    add_bullet_block(s, [
        "MTM 重构\n→ 修复 Sharpe 异常迭代 2 轮",
        "决策表 / 条件组合\n→ 需先做控制流图分析",
        "DefiLlama 时间戳归一化\n→ 跨池秒级偏移踩坑",
        "Streamlit 缓存失效\n→ 按 mtime 自动失效才解",
        "雷达图视觉缩放\n→ 拆 subplot 才让对比可读",
    ], Inches(4.9), Inches(2.5), Inches(3.75), Inches(4.3),
       font_size=11, color=C_BODY, bullet_color=C_WARN,
       line_spacing=1.4)

    # 右：降级
    add_rounded(s, Inches(8.9), Inches(1.85), Inches(4.1), Inches(5.0),
                fill=C_BG_CARD, line=C_HOT, line_width=Pt(1.2))
    add_rect(s, Inches(8.9), Inches(1.85), Inches(4.1), Inches(0.06), fill=C_HOT)
    section_title(s, Inches(9.1), Inches(2.0), "DEFERRED · 降级 / 未做", C_HOT)
    add_bullet_block(s, [
        "LP 池仅单代币建模\n（双代币 LP 留作未来）",
        "UI 仅中文\n（多语言切换未做）",
        "强化学习驱动策略\n（留作未来扩展）",
        "10000 tick 性能\n（线性外推 30 s · 接近预算）",
    ], Inches(9.1), Inches(2.5), Inches(3.75), Inches(4.3),
       font_size=11, color=C_BODY, bullet_color=C_HOT,
       line_spacing=1.5)

    add_slide_number(s, 16)
    add_speaker_notes(s,
        "整体超出预期，特别是真实数据和测试体系。也踩了一些坑，Sharpe 异常这个 bug "
        "在端到端测试后才发现，迭代了两轮才修好。"
    )

    # =====================================================
    # Slide 17：结论洞察 · KEY INSIGHTS
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_decoration_top(s)
    add_slide_title(s, "结论洞察 · KEY INSIGHTS",
                    subtitle="不仅能跑，还能回答问题 · 真实数据驱动的 3 个发现")

    # 3 个洞察大卡片，左中右排版
    insight_w = Inches(4.15)
    insight_h = Inches(5.2)
    insight_y = Inches(1.85)
    insight_gap = Inches(0.1)
    insight_x0 = Inches(0.45)

    # 洞察 ①：滑点 vs 池深（正是题目举的例子！）
    x1 = insight_x0
    add_rounded(s, x1, insight_y, insight_w, insight_h,
                fill=C_BG_CARD, line=C_ACCENT, line_width=Pt(1.2))
    add_rect(s, x1, insight_y, insight_w, Inches(0.08), fill=C_ACCENT)

    add_text(s, "Q1", x1 + Inches(0.2), insight_y + Inches(0.25),
             Inches(1.0), Inches(0.4),
             font_size=14, bold=True, color=C_ACCENT, font_name=F_MONO)
    add_text(s, "滑点如何随池深变化？",
             x1 + Inches(0.2), insight_y + Inches(0.6),
             insight_w - Inches(0.4), Inches(0.55),
             font_size=16, bold=True, color=C_TITLE, font_name=F_CN)

    # 分隔线
    add_rect(s, x1 + Inches(0.2), insight_y + Inches(1.2),
             Inches(0.4), Inches(0.025), fill=C_ACCENT)

    # 实证数据
    add_text(s, "实证 · 真实链上数据",
             x1 + Inches(0.2), insight_y + Inches(1.35),
             insight_w - Inches(0.4), Inches(0.3),
             font_size=10, color=C_DIM, font_name=F_MONO)

    # 数据对比
    pool_rows = [
        ("Maple USDC", "$3.27 B", "0.1%", C_OK),
        ("Vesper ETH", "$5.2 M",  "0.3%", C_WARN),
    ]
    py = insight_y + Inches(1.7)
    for pid, tvl, slip, col in pool_rows:
        add_rect(s, x1 + Inches(0.2), py, insight_w - Inches(0.4), Inches(0.6),
                 fill=C_BG_PANEL)
        add_text(s, pid, x1 + Inches(0.3), py,
                 Inches(1.6), Inches(0.6),
                 font_size=11, color=C_BODY, font_name=F_CN,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tvl, x1 + Inches(1.7), py,
                 Inches(1.2), Inches(0.6),
                 font_size=11, color=C_DIM, font_name=F_MONO,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_text(s, slip, x1 + insight_w - Inches(1.1), py,
                 Inches(0.9), Inches(0.6),
                 font_size=15, bold=True, color=col, font_name=F_MONO,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        py += Inches(0.7)

    # 结论 box
    add_rounded(s, x1 + Inches(0.2), insight_y + Inches(3.5),
                insight_w - Inches(0.4), Inches(1.5),
                fill=C_BG_PANEL, line=C_ACCENT, line_width=Pt(0.6))
    add_text(s, "→ 结论",
             x1 + Inches(0.35), insight_y + Inches(3.6),
             insight_w - Inches(0.7), Inches(0.35),
             font_size=11, color=C_ACCENT, bold=True, font_name=F_MONO)
    add_text(s,
             "trade / TVL 决定滑点档位\n"
             "大 TVL 池可节省 ≥ 67% 滑点\n"
             "→ 阶梯函数（FrictionEstimator）\n   合理建模这一现象",
             x1 + Inches(0.35), insight_y + Inches(3.95),
             insight_w - Inches(0.7), Inches(1.05),
             font_size=11, color=C_BODY, font_name=F_CN, line_spacing=1.4)

    # 洞察 ②：是否考虑价格风险 决定成败
    x2 = insight_x0 + insight_w + insight_gap
    add_rounded(s, x2, insight_y, insight_w, insight_h,
                fill=C_BG_CARD, line=C_HOT, line_width=Pt(1.2))
    add_rect(s, x2, insight_y, insight_w, Inches(0.08), fill=C_HOT)

    add_text(s, "Q2", x2 + Inches(0.2), insight_y + Inches(0.25),
             Inches(1.0), Inches(0.4),
             font_size=14, bold=True, color=C_HOT, font_name=F_MONO)
    add_text(s, "评分器是否考虑\ntoken_price 风险？",
             x2 + Inches(0.2), insight_y + Inches(0.6),
             insight_w - Inches(0.4), Inches(0.85),
             font_size=16, bold=True, color=C_TITLE, font_name=F_CN,
             line_spacing=1.2)

    add_rect(s, x2 + Inches(0.2), insight_y + Inches(1.5),
             Inches(0.4), Inches(0.025), fill=C_HOT)

    add_text(s, "实证 · 同数据 5 策略对比",
             x2 + Inches(0.2), insight_y + Inches(1.65),
             insight_w - Inches(0.4), Inches(0.3),
             font_size=10, color=C_DIM, font_name=F_MONO)

    # 对比卡
    py = insight_y + Inches(2.0)
    contrast = [
        ("不考虑价格", "激进动量", "-6.50%", "20.44%", C_HOT),
        ("考虑价格",   "保守稳健", "+8.82%", "5.80%",  C_OK),
    ]
    for label, name, ar, mdd, col in contrast:
        add_rect(s, x2 + Inches(0.2), py, insight_w - Inches(0.4), Inches(0.65),
                 fill=C_BG_PANEL)
        add_text(s, label, x2 + Inches(0.3), py,
                 Inches(1.5), Inches(0.3),
                 font_size=9, color=C_DIM, font_name=F_CN)
        add_text(s, name, x2 + Inches(0.3), py + Inches(0.25),
                 Inches(1.5), Inches(0.4),
                 font_size=11, bold=True, color=col, font_name=F_CN)
        add_text(s, ar, x2 + Inches(1.85), py + Inches(0.05),
                 Inches(1.2), Inches(0.3),
                 font_size=12, bold=True, color=col, font_name=F_MONO,
                 align=PP_ALIGN.CENTER)
        add_text(s, "年化", x2 + Inches(1.85), py + Inches(0.4),
                 Inches(1.2), Inches(0.25),
                 font_size=8, color=C_DIM, font_name=F_CN,
                 align=PP_ALIGN.CENTER)
        add_text(s, mdd, x2 + insight_w - Inches(1.2), py + Inches(0.05),
                 Inches(1.0), Inches(0.3),
                 font_size=12, bold=True, color=col, font_name=F_MONO,
                 align=PP_ALIGN.CENTER)
        add_text(s, "MDD", x2 + insight_w - Inches(1.2), py + Inches(0.4),
                 Inches(1.0), Inches(0.25),
                 font_size=8, color=C_DIM, font_name=F_CN,
                 align=PP_ALIGN.CENTER)
        py += Inches(0.72)

    # 结论
    add_rounded(s, x2 + Inches(0.2), insight_y + Inches(3.55),
                insight_w - Inches(0.4), Inches(1.45),
                fill=C_BG_PANEL, line=C_HOT, line_width=Pt(0.6))
    add_text(s, "→ 结论",
             x2 + Inches(0.35), insight_y + Inches(3.65),
             insight_w - Inches(0.7), Inches(0.35),
             font_size=11, color=C_HOT, bold=True, font_name=F_MONO)
    add_text(s,
             "光看 APY 是不够的\n"
             "必须把 token 价格风险纳入评分\n"
             "→ TokenPriceVol/MDDPenalty\n   解决「-6.5% 亏损」问题",
             x2 + Inches(0.35), insight_y + Inches(4.0),
             insight_w - Inches(0.7), Inches(1.0),
             font_size=11, color=C_BODY, font_name=F_CN, line_spacing=1.4)

    # 洞察 ③：调仓数量 vs 收益
    x3 = insight_x0 + 2 * (insight_w + insight_gap)
    add_rounded(s, x3, insight_y, insight_w, insight_h,
                fill=C_BG_CARD, line=C_OK, line_width=Pt(1.2))
    add_rect(s, x3, insight_y, insight_w, Inches(0.08), fill=C_OK)

    add_text(s, "Q3", x3 + Inches(0.2), insight_y + Inches(0.25),
             Inches(1.0), Inches(0.4),
             font_size=14, bold=True, color=C_OK, font_name=F_MONO)
    add_text(s, "调仓越频繁\n收益越高吗？",
             x3 + Inches(0.2), insight_y + Inches(0.6),
             insight_w - Inches(0.4), Inches(0.85),
             font_size=16, bold=True, color=C_TITLE, font_name=F_CN,
             line_spacing=1.2)

    add_rect(s, x3 + Inches(0.2), insight_y + Inches(1.5),
             Inches(0.4), Inches(0.025), fill=C_OK)

    add_text(s, "实证 · 5 预设统计",
             x3 + Inches(0.2), insight_y + Inches(1.65),
             insight_w - Inches(0.4), Inches(0.3),
             font_size=10, color=C_DIM, font_name=F_MONO)

    # 散点对比
    rotation_data = [
        ("激进",     "1 次",  "-6.5%",  C_HOT),
        ("低频",     "16 次", "-11.7%", C_HOT),
        ("保守",     "19 次", "+8.8%",  C_OK),
    ]
    py = insight_y + Inches(2.0)
    for name, n, ar, col in rotation_data:
        add_rect(s, x3 + Inches(0.2), py, insight_w - Inches(0.4), Inches(0.42),
                 fill=C_BG_PANEL)
        add_text(s, name, x3 + Inches(0.35), py,
                 Inches(1.2), Inches(0.42),
                 font_size=11, bold=True, color=C_BODY, font_name=F_CN,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, n, x3 + Inches(1.6), py,
                 Inches(1.2), Inches(0.42),
                 font_size=11, color=C_DIM, font_name=F_MONO,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        add_text(s, ar, x3 + insight_w - Inches(1.2), py,
                 Inches(1.0), Inches(0.42),
                 font_size=13, bold=True, color=col, font_name=F_MONO,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        py += Inches(0.5)

    add_rounded(s, x3 + Inches(0.2), insight_y + Inches(3.55),
                insight_w - Inches(0.4), Inches(1.45),
                fill=C_BG_PANEL, line=C_OK, line_width=Pt(0.6))
    add_text(s, "→ 结论",
             x3 + Inches(0.35), insight_y + Inches(3.65),
             insight_w - Inches(0.7), Inches(0.35),
             font_size=11, color=C_OK, bold=True, font_name=F_MONO)
    add_text(s,
             "调仓数量不是关键\n"
             "调仓时机才是关键\n"
             "→ τ-reset + 双门槛 + 价格漂移\n   联合决定何时该出手",
             x3 + Inches(0.35), insight_y + Inches(4.0),
             insight_w - Inches(0.7), Inches(1.0),
             font_size=11, color=C_BODY, font_name=F_CN, line_spacing=1.4)

    # 底部总结
    add_rounded(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
                fill=C_BG_CARD, line=C_LINE, line_width=Pt(0.5))
    add_text(s,
             "回答问题型回测平台  ·  数据驱动结论  ·  非「跑通即止」",
             Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
             font_size=11, color=C_DIM, italic=True, font_name=F_CN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_slide_number(s, 17)
    add_speaker_notes(s,
        "三个洞察都基于真实数据。"
        "Q1：滑点和池深的关系——大 TVL 池可省 67% 滑点；"
        "Q2：评分器是否考虑价格风险决定盈亏，是赚 8.8% 还是亏 6.5% 的分水岭；"
        "Q3：调仓数量不是关键，时机才是——激进 1 次亏、低频 16 次亏、保守 19 次赚。"
        "这三个结论都是平台跑出来的，不是我们事先知道的。"
    )

    # =====================================================
    # Slide 18：总结 + Q&A
    # =====================================================
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    add_rect(s, Emu(0), Emu(0), SLIDE_W, Inches(0.08), fill=C_ACCENT)
    add_rect(s, Emu(0), Inches(0.08), SLIDE_W, Inches(0.02), fill=C_HOT)

    # 大字标语
    add_text(s, "完整可运行 · 可复现 · 可扩展",
             Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.0),
             font_size=42, bold=True, color=C_TITLE, font_name=F_CN,
             align=PP_ALIGN.CENTER)
    add_text(s, "RUNNABLE  ·  REPRODUCIBLE  ·  EXTENSIBLE",
             Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.5),
             font_size=14, color=C_ACCENT, font_name=F_MONO,
             align=PP_ALIGN.CENTER)
    add_rect(s, Inches(6.0), Inches(2.95), Inches(1.3), Inches(0.04), fill=C_HOT)

    # 关键数字回顾（7 个圆角卡）
    keys = [
        ("13.5k+", "代码行"),
        ("325",    "测试"),
        ("87.5%",  "覆盖率"),
        ("6",      "Tab"),
        ("10",     "图表"),
        ("5",      "预设"),
        ("REAL",   "链上数据"),
    ]
    metric_w = Inches(1.6)
    total_w = len(keys) * metric_w + (len(keys) - 1) * Inches(0.12)
    metric_x0 = (SLIDE_W - total_w) / 2
    metric_y = Inches(3.5)
    for i, (num, lab) in enumerate(keys):
        x = metric_x0 + i * (metric_w + Inches(0.12))
        add_rounded(s, x, metric_y, metric_w, Inches(1.5),
                    fill=C_BG_CARD, line=C_ACCENT, line_width=Pt(0.8))
        add_text(s, num, x, metric_y + Inches(0.15),
                 metric_w, Inches(0.7),
                 font_size=20, bold=True, color=C_ACCENT,
                 font_name=F_MONO, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, lab, x, metric_y + Inches(0.9),
                 metric_w, Inches(0.5),
                 font_size=10, color=C_DIM,
                 font_name=F_CN, align=PP_ALIGN.CENTER)

    # Q & A
    add_text(s, "Q  &  A",
             Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.8),
             font_size=40, bold=True, color=C_HOT, font_name=F_CN,
             align=PP_ALIGN.CENTER)

    # 命令小字
    add_text(s,
             "$ pytest tests/  ·  python run_example.py  ·  streamlit run ui/app.py",
             Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
             font_size=11, color=C_DIM, font_name=F_MONO,
             align=PP_ALIGN.CENTER)

    add_text(s, "谢 谢 聆 听   ·   欢 迎 提 问",
             Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
             font_size=14, color=C_BODY, font_name=F_CN,
             align=PP_ALIGN.CENTER)

    add_speaker_notes(s, "总结：完整可运行、可复现、可扩展。谢谢大家，欢迎提问。")

    # 保存
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    size_kb = OUT.stat().st_size / 1024
    print(f"[saved] {OUT.relative_to(ROOT)}  ({size_kb:.0f} KB)")
    print(f"幻灯片数：{len(prs.slides)}")


if __name__ == "__main__":
    main()
